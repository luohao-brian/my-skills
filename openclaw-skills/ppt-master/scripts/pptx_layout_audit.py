#!/usr/bin/env python3
"""Audit geometry in an exported PPTX package.

The SVG browser audit checks the authored design source. This script provides a
second, export-side check using the actual DrawingML geometry exposed by
python-pptx.

Exit codes:
    0  no blocking errors
    1  blocking geometry errors (or warnings with --fail-on-warnings)
    2  invalid input/dependency/package
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any, Iterable


def _rect(shape: Any) -> dict[str, int]:
    return {
        "x": int(shape.left),
        "y": int(shape.top),
        "width": int(shape.width),
        "height": int(shape.height),
    }


def _area(rect: dict[str, int]) -> int:
    return max(0, rect["width"]) * max(0, rect["height"])


def _intersection(a: dict[str, int], b: dict[str, int]) -> dict[str, int]:
    x1, y1 = max(a["x"], b["x"]), max(a["y"], b["y"])
    x2 = min(a["x"] + a["width"], b["x"] + b["width"])
    y2 = min(a["y"] + a["height"], b["y"] + b["height"])
    return {"x": x1, "y": y1, "width": max(0, x2 - x1), "height": max(0, y2 - y1)}


def _contains(outer: dict[str, int], inner: dict[str, int], tolerance: int) -> bool:
    return (
        inner["x"] >= outer["x"] - tolerance
        and inner["y"] >= outer["y"] - tolerance
        and inner["x"] + inner["width"] <= outer["x"] + outer["width"] + tolerance
        and inner["y"] + inner["height"] <= outer["y"] + outer["height"] + tolerance
    )


def _iter_leaf_shapes(shapes: Iterable[Any], parent: str = "slide") -> Iterable[tuple[Any, str]]:
    for index, shape in enumerate(shapes, start=1):
        name = str(getattr(shape, "name", "") or f"shape-{index}")
        path = f"{parent}/{name}"
        if hasattr(shape, "shapes"):
            yield from _iter_leaf_shapes(shape.shapes, path)
        else:
            yield shape, path


def _shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    try:
        return str(shape.text or "").strip()
    except Exception:  # noqa: BLE001
        return ""


def _union(rectangles: list[dict[str, int]]) -> dict[str, int] | None:
    if not rectangles:
        return None
    x1 = min(rect["x"] for rect in rectangles)
    y1 = min(rect["y"] for rect in rectangles)
    x2 = max(rect["x"] + rect["width"] for rect in rectangles)
    y2 = max(rect["y"] + rect["height"] for rect in rectangles)
    return {"x": x1, "y": y1, "width": x2 - x1, "height": y2 - y1}


def audit_pptx(
    pptx_path: Path,
    min_width_use: float,
    min_height_use: float,
) -> dict[str, Any]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required; install requirements.txt") from exc

    presentation = Presentation(str(pptx_path))
    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)
    slide_box = {"x": 0, "y": 0, "width": slide_width, "height": slide_height}
    tolerance = max(1, int(min(slide_width, slide_height) * 0.001))
    pages: list[dict[str, Any]] = []
    issues: list[dict[str, Any]] = []

    for page_number, slide in enumerate(presentation.slides, start=1):
        records: list[dict[str, Any]] = []
        for shape, path in _iter_leaf_shapes(slide.shapes):
            box = _rect(shape)
            text = _shape_text(shape)
            record = {
                "name": path,
                "parent": path.rsplit("/", 1)[0],
                "bbox": box,
                "text": text,
            }
            records.append(record)
            if box["width"] <= 0 or box["height"] <= 0:
                issues.append(
                    {
                        "severity": "error",
                        "code": "zero_size_shape",
                        "slide": page_number,
                        "shape": path,
                        "detail": f"non-positive shape bounds: {box}",
                    }
                )
            elif not _contains(slide_box, box, tolerance):
                issues.append(
                    {
                        "severity": "error",
                        "code": "shape_outside_slide",
                        "slide": page_number,
                        "shape": path,
                        "detail": f"shape bounds {box} exceed slide bounds {slide_box}",
                    }
                )

        text_records = [record for record in records if record["text"]]
        for index, left in enumerate(text_records):
            for right in text_records[index + 1 :]:
                hit = _intersection(left["bbox"], right["bbox"])
                if hit["width"] <= tolerance or hit["height"] <= tolerance:
                    continue
                ratio = _area(hit) / max(1, min(_area(left["bbox"]), _area(right["bbox"])))
                same_parent = left["parent"] == right["parent"]
                if same_parent and ratio <= 0.75:
                    continue
                if not same_parent and ratio <= 0.12:
                    continue
                severity = "error" if ratio > 0.75 else "warning"
                issues.append(
                    {
                        "severity": severity,
                        "code": "text_frame_overlap",
                        "slide": page_number,
                        "shape": left["name"],
                        "other": right["name"],
                        "detail": f"text frames overlap; intersection ratio {ratio:.3f}",
                    }
                )

        meaningful: list[dict[str, int]] = []
        for record in records:
            box = record["bbox"]
            width_ratio = box["width"] / slide_width
            height_ratio = box["height"] / slide_height
            if width_ratio > 0.92 and height_ratio > 0.92 and not record["text"]:
                continue
            if box["width"] <= tolerance or box["height"] <= tolerance:
                continue
            meaningful.append(box)
        content_box = _union(meaningful)
        if content_box is None:
            utilization = {
                "width": 0.0,
                "height": 0.0,
                "left": 1.0,
                "right": 1.0,
                "top": 1.0,
                "bottom": 1.0,
            }
        else:
            utilization = {
                "width": content_box["width"] / slide_width,
                "height": content_box["height"] / slide_height,
                "left": content_box["x"] / slide_width,
                "right": (slide_width - content_box["x"] - content_box["width"]) / slide_width,
                "top": content_box["y"] / slide_height,
                "bottom": (slide_height - content_box["y"] - content_box["height"]) / slide_height,
            }
        if utilization["width"] < min_width_use or utilization["height"] < min_height_use:
            issues.append(
                {
                    "severity": "warning",
                    "code": "underused_slide",
                    "slide": page_number,
                    "shape": None,
                    "detail": (
                        f"content span {utilization['width'] * 100:.1f}% × "
                        f"{utilization['height'] * 100:.1f}%"
                    ),
                    "utilization": utilization,
                }
            )
        tight_sides = [
            side for side in ("left", "right", "top", "bottom")
            if utilization[side] < 0.008
        ]
        if len(tight_sides) >= 3:
            issues.append(
                {
                    "severity": "warning",
                    "code": "edge_crowding",
                    "slide": page_number,
                    "shape": None,
                    "detail": f"content is tight to {', '.join(tight_sides)}",
                    "utilization": utilization,
                }
            )
        pages.append(
            {
                "slide": page_number,
                "shape_count": len(records),
                "text_shape_count": len(text_records),
                "content_bbox": content_box,
                "utilization": utilization,
            }
        )

    errors = [issue for issue in issues if issue["severity"] == "error"]
    warnings = [issue for issue in issues if issue["severity"] == "warning"]
    return {
        "pptx": str(pptx_path),
        "slide_size_emu": {"width": slide_width, "height": slide_height},
        "slide_count": len(presentation.slides),
        "summary": {"errors": len(errors), "warnings": len(warnings)},
        "issues": issues,
        "pages": pages,
    }


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Audit exported PPTX slide geometry.")
    parser.add_argument("pptx_path", help="Exported .pptx file")
    parser.add_argument("--report", help="JSON report path (default: beside PPTX)")
    parser.add_argument("--min-width-use", type=float, default=0.48)
    parser.add_argument("--min-height-use", type=float, default=0.38)
    parser.add_argument("--fail-on-warnings", action="store_true")
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    pptx_path = Path(args.pptx_path).expanduser().resolve()
    if not pptx_path.is_file() or pptx_path.suffix.lower() != ".pptx":
        print(f"PPTX file not found: {pptx_path}", file=sys.stderr)
        return 2
    if not (0 < args.min_width_use <= 1 and 0 < args.min_height_use <= 1):
        print("utilization thresholds must be in (0, 1]", file=sys.stderr)
        return 2
    try:
        report = audit_pptx(pptx_path, args.min_width_use, args.min_height_use)
    except Exception as exc:  # noqa: BLE001
        print(f"PPTX audit failed: {type(exc).__name__}: {exc}", file=sys.stderr)
        return 2
    report_path = (
        Path(args.report).expanduser().resolve()
        if args.report
        else pptx_path.with_suffix(".layout-audit.json")
    )
    report_path.parent.mkdir(parents=True, exist_ok=True)
    report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    summary = report["summary"]
    print(
        f"PPTX audit: {report['slide_count']} slide(s), "
        f"{summary['errors']} error(s), {summary['warnings']} warning(s)"
    )
    for issue in report["issues"][:40]:
        print(
            f"[{issue['severity'].upper()}] slide {issue['slide']} {issue['code']} "
            f"{issue.get('shape') or ''}: {issue['detail']}"
        )
    if len(report["issues"]) > 40:
        print(f"... {len(report['issues']) - 40} more issue(s); see report")
    print(f"Report: {report_path}")
    if summary["errors"] or (args.fail_on_warnings and summary["warnings"]):
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
